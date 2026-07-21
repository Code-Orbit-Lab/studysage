"""PPTX parsing — python-pptx. Owner: Saurabh"""
from pathlib import Path
from pptx import Presentation


def parse_pptx(path: Path) -> list[dict]:
    """Returns [{"page": int, "text": str}, ...] — "page" here means slide number."""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        combined = "\n".join(t for t in texts if t.strip())
        if combined.strip():
            slides.append({"page": i, "text": combined})
    return slides